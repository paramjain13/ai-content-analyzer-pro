from pptx import Presentation
import re

def extract_text_from_pptx(file_path):
    """
    Extract text from PPTX file
    Returns: (text, metadata)
    """
    try:
        prs = Presentation(file_path)
        
        # Extract metadata
        metadata = {
            "title": prs.core_properties.title or "Untitled Presentation",
            "author": prs.core_properties.author or "Unknown",
            "created": str(prs.core_properties.created) if prs.core_properties.created else "Unknown",
            "modified": str(prs.core_properties.modified) if prs.core_properties.modified else "Unknown",
            "slides": len(prs.slides)
        }
        
        # Extract text from all slides
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            # Add slide marker
            text_parts.append(f"[Slide {slide_num}]")
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        text_parts.append(text)
                
                # Extract text from tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                text_parts.append(cell_text)
        
        full_text = " ".join(text_parts)
        
        # Clean up whitespace
        full_text = re.sub(r'\s+', ' ', full_text).strip()
        
        if not full_text or len(full_text) < 50:
            raise Exception("No readable text found in PPTX file")
        
        return full_text, metadata
        
    except Exception as e:
        raise Exception(f"Error reading PPTX file: {str(e)}")